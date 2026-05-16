#!/usr/bin/env python3
"""
scout_corpus.py — фаза 2 пайплайна doc2kb. Сканирует корпус и составляет
карту для последующих фаз.

CLI:
    scout_corpus.py <input_dir> <kb_dir>

Эффекты:
    - Создаёт <kb_dir>/_scout.json по схеме references/format-spec.md.
    - Пишет короткое summary на stdout (по 1 строке на тип файла).

Не делает:
    - Не извлекает контент (это фаза 4).
    - Не спрашивает у пользователя — список решений собирается в
      _scout.json.user_decisions_needed; диалог с пользователем ведёт
      главный агент по инструкции из SKILL.md.

Принципы:
    - Никогда не падать целиком из-за одного файла. Любая ошибка по
      конкретному файлу → запись в files[] с warnings + skipped_at_scout
      при невозможности классифицировать.
    - Никогда не доверять расширению целиком: mime определяется через
      libmagic (python-magic), сравнивается с extension; при расхождении
      mime_confidence: "low".
    - Тихо игнорировать Office lock-files (~$*.docx, ~$*.pptx), .DS_Store,
      Thumbs.db, скрытые директории (.git, node_modules и т.п.).
"""
from __future__ import annotations

import argparse
import hashlib
import json
import re
import sys
import time
import unicodedata
from datetime import datetime, timezone
from pathlib import Path
from typing import Any


# Schema version of the scout JSON output. Bump on breaking changes.
SCHEMA_VERSION = "1.0"


def _scout_tool_id() -> str:
    """Lazy because tool_version_string walks parents — keep import-time cheap."""
    try:
        # _common is alongside this script.
        from _common import tool_version_string  # type: ignore
        return f"doc2kb@{tool_version_string()}"
    except Exception:
        return "doc2kb@unknown"
# Minimum text length per page to treat a PDF page as having an embedded text layer.
PDF_TEXT_CHAR_THRESHOLD = 100
# Tunables for "huge file" warnings.
HUGE_BYTES = 50 * 1024 * 1024      # 50 MB
HUGE_PAGES = 500                    # PDF pages

# Lock-files and OS noise — silent skip, no entry in skipped_at_scout.
SILENT_SKIP_NAMES = {".DS_Store", "Thumbs.db", "desktop.ini"}
SILENT_SKIP_PREFIXES = ("~$", "._")   # Office lock-files and macOS metadata
SILENT_SKIP_DIRS = {
    ".git", ".hg", ".svn", "node_modules", "__pycache__",
    ".venv", "venv", ".tox", ".mypy_cache", ".pytest_cache",
    ".idea", ".vscode", ".cache",
}

# Mapping ext → source_type. Lowercase. Extensions outside this map go through
# mime-based detection and may end up as "unknown".
EXT_TO_TYPE = {
    ".pdf": "pdf",
    ".docx": "docx",
    ".pptx": "pptx",
    ".xlsx": "xlsx",
    ".md": "md",
    ".markdown": "md",
    ".txt": "txt",
    ".html": "html",
    ".htm": "html",
    ".epub": "epub",
    ".rtf": "rtf",
    ".odt": "odt",
    ".png": "image",
    ".jpg": "image",
    ".jpeg": "image",
    ".gif": "image",
    ".tiff": "image",
    ".tif": "image",
    ".webp": "image",
}

# Source-types supported by MVP (have a dedicated extract_*.py).
MVP_SUPPORTED = {"pdf", "docx", "pptx", "md", "txt", "html"}


# ---------- low-level helpers ----------

def log(msg: str) -> None:
    print(f"[doc2kb scout] {msg}", file=sys.stderr, flush=True)


def sha256_of(path: Path, chunk: int = 1 << 20) -> str:
    h = hashlib.sha256()
    with path.open("rb") as fh:
        while True:
            buf = fh.read(chunk)
            if not buf:
                break
            h.update(buf)
    return h.hexdigest()


def slugify(text: str, maxlen: int = 48) -> str:
    """ASCII-friendly slug. Preserves Cyrillic transliteration roughly through
    NFKD when possible; otherwise drops non-ASCII. Result lowercased,
    non-alnum→'-', deduplicated, trimmed."""
    s = unicodedata.normalize("NFKD", text)
    s = s.encode("ascii", "ignore").decode("ascii")
    s = re.sub(r"[^A-Za-z0-9]+", "-", s).strip("-").lower()
    if not s:
        s = "doc"
    return s[:maxlen].rstrip("-") or "doc"


def sniff_mime(path: Path) -> str | None:
    """Returns mime via libmagic or None if libmagic is unavailable/failed."""
    try:
        import magic  # type: ignore
    except Exception:
        return None
    try:
        return magic.from_file(str(path), mime=True)
    except Exception:
        return None


def detect_source_type(path: Path, mime: str | None) -> tuple[str, str]:
    """Returns (source_type, confidence). Compares extension and mime.

    confidence:
      "high" — extension and mime agree (or only one is conclusive).
      "low"  — extension says one thing, mime says another.
    """
    ext = path.suffix.lower()
    ext_type = EXT_TO_TYPE.get(ext, "unknown")

    mime_type = "unknown"
    if mime:
        if mime == "application/pdf":
            mime_type = "pdf"
        elif mime in (
            "application/vnd.openxmlformats-officedocument.wordprocessingml.document",
        ):
            mime_type = "docx"
        elif mime in (
            "application/vnd.openxmlformats-officedocument.presentationml.presentation",
        ):
            mime_type = "pptx"
        elif mime in (
            "application/vnd.openxmlformats-officedocument.spreadsheetml.sheet",
        ):
            mime_type = "xlsx"
        elif mime in ("text/markdown",):
            mime_type = "md"
        elif mime.startswith("text/html"):
            mime_type = "html"
        elif mime in ("application/epub+zip",):
            mime_type = "epub"
        elif mime in ("application/rtf", "text/rtf"):
            mime_type = "rtf"
        elif mime in ("application/vnd.oasis.opendocument.text",):
            mime_type = "odt"
        elif mime.startswith("image/"):
            mime_type = "image"
        elif mime.startswith("text/"):
            # generic text — treat as txt, but only if extension agrees
            mime_type = "txt"

    if ext_type == "unknown" and mime_type != "unknown":
        return mime_type, "high"
    if ext_type != "unknown" and mime_type == "unknown":
        return ext_type, "high"
    if ext_type == mime_type:
        return ext_type, "high"
    # Disagreement — prefer mime (less spoofable), flag low confidence.
    return mime_type, "low"


# ---------- per-type classifiers ----------

def classify_pdf(path: Path) -> dict[str, Any]:
    """Returns subset of file_info — pdf_class, pages, encryption flag."""
    out: dict[str, Any] = {
        "pdf_class": "unknown",
        "pages": None,
        "has_tables": False,
        "has_images": False,
    }
    # Encryption check first — pdfplumber will fail on encrypted PDFs.
    try:
        import pikepdf  # type: ignore
    except Exception:
        pikepdf = None  # type: ignore
    if pikepdf is not None:
        try:
            with pikepdf.open(str(path)) as pdf:
                out["pages"] = len(pdf.pages)
        except pikepdf.PasswordError:  # type: ignore
            out["pdf_class"] = "encrypted"
            return out
        except Exception:
            # Fall through to pdfplumber for page count + classification.
            # The pikepdf-specific error doesn't carry forward — pdfplumber
            # will either succeed (recovering page count) or set
            # pdf_class="corrupt" with its own warning.
            pass

    # Layer detection via pdfplumber.
    try:
        import pdfplumber  # type: ignore
    except Exception:
        out["pdf_class"] = "unknown"
        out["warnings"] = ["pdfplumber unavailable"]
        return out
    try:
        with pdfplumber.open(str(path)) as pdf:
            n_pages = len(pdf.pages)
            if out["pages"] is None:
                out["pages"] = n_pages
            text_pages = 0
            image_pages = 0
            table_pages = 0
            # Sample up to 12 pages evenly across the document — full scan on
            # 1000-page books is too slow for scout phase.
            if n_pages <= 12:
                indices = list(range(n_pages))
            else:
                indices = sorted({
                    int(i * (n_pages - 1) / 11) for i in range(12)
                })
            sampled = len(indices)
            for idx in indices:
                page = pdf.pages[idx]
                try:
                    txt = page.extract_text() or ""
                except Exception:
                    txt = ""
                if len(txt.strip()) >= PDF_TEXT_CHAR_THRESHOLD:
                    text_pages += 1
                try:
                    if page.images:
                        image_pages += 1
                except Exception:
                    pass
                try:
                    if page.find_tables():
                        table_pages += 1
                except Exception:
                    pass

            if text_pages == 0:
                out["pdf_class"] = "image_only"
            elif text_pages < sampled * 0.5:
                out["pdf_class"] = "mixed"
            else:
                out["pdf_class"] = "text"
            out["has_tables"] = table_pages > 0
            out["has_images"] = image_pages > 0
    except Exception as e:
        # Could be password-protected without raising in pikepdf (rare),
        # or genuinely corrupt.
        msg = str(e).lower()
        if "password" in msg or "encrypt" in msg:
            out["pdf_class"] = "encrypted"
        else:
            out["pdf_class"] = "corrupt"
            out["warnings"] = [f"pdf classification failed: {str(e)[:200]}"]
    return out


def docx_richness(path: Path) -> dict[str, Any]:
    out: dict[str, Any] = {
        "inline_images": 0,
        "has_tables": False,
        "has_equations": False,
        "has_tracked_changes": False,
        "paragraphs": 0,
    }
    try:
        from docx import Document  # type: ignore
    except Exception:
        out["warnings"] = ["python-docx unavailable"]
        return out
    try:
        doc = Document(str(path))
        out["paragraphs"] = len(doc.paragraphs)
        out["inline_images"] = len(doc.inline_shapes)
        out["has_tables"] = len(doc.tables) > 0
        xml = doc.element.xml
        out["has_equations"] = "<m:oMath" in xml or "<m:oMathPara" in xml
        out["has_tracked_changes"] = "w:ins" in xml or "w:del" in xml
    except Exception as e:
        out["warnings"] = [f"docx scout failed: {str(e)[:200]}"]
    return out


def pptx_richness(path: Path) -> dict[str, Any]:
    out: dict[str, Any] = {
        "slides": 0,
        "has_notes": False,
        "notes_chars": 0,
        "inline_images": 0,
        "has_tables": False,
    }
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
    except Exception:
        out["warnings"] = ["python-pptx unavailable"]
        return out
    try:
        prs = Presentation(str(path))
        out["slides"] = len(prs.slides)
        n_images = 0
        n_tables = 0
        notes_chars = 0
        any_notes = False
        for slide in prs.slides:
            if slide.has_notes_slide:
                try:
                    t = slide.notes_slide.notes_text_frame.text or ""
                except Exception:
                    t = ""
                t = t.strip()
                if t:
                    any_notes = True
                    notes_chars += len(t)
            for sh in slide.shapes:
                try:
                    if sh.shape_type == MSO_SHAPE_TYPE.PICTURE:
                        n_images += 1
                    if getattr(sh, "has_table", False):
                        n_tables += 1
                except Exception:
                    pass
        out["has_notes"] = any_notes
        out["notes_chars"] = notes_chars
        out["inline_images"] = n_images
        out["has_tables"] = n_tables > 0
    except Exception as e:
        out["warnings"] = [f"pptx scout failed: {str(e)[:200]}"]
    return out


def text_encoding(path: Path, sample_bytes: int = 65536) -> dict[str, Any]:
    out: dict[str, Any] = {"encoding": "utf-8"}
    try:
        from charset_normalizer import from_bytes  # type: ignore
    except Exception:
        return out
    try:
        with path.open("rb") as fh:
            data = fh.read(sample_bytes)
        result = from_bytes(data).best()
        if result is not None:
            out["encoding"] = result.encoding or "utf-8"
    except Exception:
        pass
    return out


# ---------- strategy + token estimate ----------

def set_strategy_and_tokens(info: dict[str, Any]) -> None:
    """Sets extraction_strategy and estimated_tokens on info in place."""
    t = info.get("source_type", "unknown")
    warnings: list[str] = info.get("warnings", []) or []

    if t == "pdf":
        pdf_class = info.get("pdf_class")
        pages = info.get("pages") or 0
        if pdf_class == "encrypted":
            info["extraction_strategy"] = "needs_password"
            info["estimated_tokens"] = None
            info["action_required"] = "ask_user_password_or_skip"
        elif pdf_class == "image_only":
            info["extraction_strategy"] = "needs_ocr_or_vlm"
            info["estimated_tokens"] = None
            info["action_required"] = "ask_user_ocr_strategy"
        elif pdf_class == "mixed":
            info["extraction_strategy"] = "pymupdf4llm"
            info["estimated_tokens"] = pages * 500 if pages else None
            warnings.append("mixed pdf — some pages may be image-only")
        elif pdf_class == "corrupt":
            info["extraction_strategy"] = "skip"
            info["estimated_tokens"] = None
            info["action_required"] = "ask_user_skip_corrupt"
        else:  # text
            info["extraction_strategy"] = "pymupdf4llm"
            info["estimated_tokens"] = pages * 600 if pages else None
        if pages and pages > HUGE_PAGES:
            warnings.append(f"very large pdf: {pages} pages")
            info["action_required"] = info.get("action_required") or "ask_user_proceed_huge"
    elif t == "docx":
        # Cheap heuristic: paragraphs * ~25 tokens average. Will be refined
        # after extraction by token_count.py.
        para = info.get("paragraphs") or 0
        info["extraction_strategy"] = "mammoth"
        info["estimated_tokens"] = max(para * 25, 200) if para else None
    elif t == "pptx":
        slides = info.get("slides") or 0
        notes = info.get("notes_chars") or 0
        # slides: ~200 tok per slide body; notes: ~1 tok per 3 chars (russian)
        info["extraction_strategy"] = "python-pptx"
        info["estimated_tokens"] = slides * 200 + notes // 3 if slides else None
    elif t == "md":
        info["extraction_strategy"] = "passthrough-md"
        info["estimated_tokens"] = max(info.get("size_bytes", 0) // 4, 100)
    elif t == "txt":
        info["extraction_strategy"] = "passthrough-txt"
        info["estimated_tokens"] = max(info.get("size_bytes", 0) // 4, 100)
    elif t == "html":
        info["extraction_strategy"] = "trafilatura"
        info["estimated_tokens"] = max(info.get("size_bytes", 0) // 6, 100)
    elif t in {"xlsx", "epub", "rtf", "odt", "image"}:
        info["extraction_strategy"] = "not_in_mvp"
        info["estimated_tokens"] = None
        warnings.append(f"{t}: not supported in MVP (follow-up release)")
        info["action_required"] = "ask_user_skip_unsupported"
    else:
        info["extraction_strategy"] = "skip"
        info["estimated_tokens"] = None
        warnings.append("unknown source type")

    # Size-based escalation applies to any extractable strategy, not just PDF.
    # A 60 MB DOCX or 80 MB PPTX deserves the same user-confirmation gate as
    # a 1000-page PDF.
    if info.get("size_bytes", 0) > HUGE_BYTES:
        warnings.append(f"large file: {info['size_bytes'] // (1024 * 1024)} MB")
        if info.get("extraction_strategy") in (
            "pymupdf4llm", "mammoth", "python-pptx",
            "passthrough-md", "passthrough-txt", "trafilatura",
        ) and not info.get("action_required"):
            info["action_required"] = "ask_user_proceed_huge"

    info["warnings"] = warnings


# ---------- walker + main ----------

def is_skipped_dir(p: Path) -> bool:
    return p.name in SILENT_SKIP_DIRS or p.name.startswith(".")


def is_silent_skip(p: Path) -> bool:
    if p.name in SILENT_SKIP_NAMES:
        return True
    for prefix in SILENT_SKIP_PREFIXES:
        if p.name.startswith(prefix):
            return True
    return False


def _resolves_inside(path: Path, real_root: Path) -> bool:
    """Returns True iff path.resolve() is real_root or a descendant of it.
    Catches symlinks pointing outside the corpus root, which would otherwise
    let a hostile corpus exfiltrate /etc/passwd or ~/.ssh/* into the kb."""
    try:
        resolved = path.resolve()
    except (OSError, RuntimeError):
        return False
    if resolved == real_root:
        return True
    return real_root in resolved.parents


def walk_corpus(root: Path) -> tuple[list[Path], list[dict]]:
    """Returns (paths, escapes). `escapes` are symlink/junction entries that
    resolve outside the corpus root — they go into skipped_at_scout with a
    clear reason so the operator sees the rejected files."""
    real_root = root.resolve()
    out: list[Path] = []
    escapes: list[dict] = []
    # os.walk with followlinks=False is the only stdlib way to traverse a
    # tree without descending into symlinked directories (Path.rglob follows
    # them pre-3.13). Then we explicitly reject symlink files whose target
    # is outside the corpus root.
    import os as _os
    for dirpath, dirnames, filenames in _os.walk(root, followlinks=False):
        # Prune dot-dirs and known skip-dirs in place.
        dirnames[:] = [d for d in dirnames if not (
            d in SILENT_SKIP_DIRS or d.startswith(".")
        )]
        for fname in sorted(filenames):
            path = Path(dirpath) / fname
            if is_silent_skip(path):
                continue
            if path.is_symlink():
                # Honour symlinks only if they stay inside the corpus.
                if not _resolves_inside(path, real_root):
                    try:
                        rel = path.relative_to(root)
                    except ValueError:
                        rel = path
                    escapes.append({
                        "source_path": str(rel),
                        "reason": "symlink escapes corpus root — refused (security)",
                    })
                    continue
            if not path.is_file():
                continue
            out.append(path)
    out.sort()
    return out, escapes


def scan_file(idx: int, path: Path, input_root: Path) -> dict[str, Any]:
    rel = path.relative_to(input_root)
    info: dict[str, Any] = {
        "id": f"doc-{idx:03d}",
        "source_path": str(rel),
        "size_bytes": path.stat().st_size,
        "sha256": sha256_of(path),
        "mime": None,
        "mime_confidence": "high",
        "source_type": "unknown",
        "warnings": [],
        "action_required": None,
    }
    mime = sniff_mime(path)
    info["mime"] = mime
    src_type, conf = detect_source_type(path, mime)
    info["source_type"] = src_type
    info["mime_confidence"] = conf
    if conf == "low":
        info["warnings"].append(f"mime/extension mismatch (mime={mime}, ext={path.suffix})")

    # Per-type augmentation.
    if src_type == "pdf":
        info.update(classify_pdf(path))
    elif src_type == "docx":
        info.update(docx_richness(path))
    elif src_type == "pptx":
        info.update(pptx_richness(path))
    elif src_type in ("md", "txt"):
        info.update(text_encoding(path))

    set_strategy_and_tokens(info)

    # Normalize: ensure known keys present even when None — keeps schema stable.
    info.setdefault("pages", None)
    info.setdefault("slides", None)
    info.setdefault("has_notes", None)
    info.setdefault("notes_chars", 0)
    info.setdefault("inline_images", 0)
    info.setdefault("has_tables", False)
    info.setdefault("has_equations", False)
    info.setdefault("encoding", None)
    info.setdefault("pdf_class", None)
    return info


def aggregate_decisions(files: list[dict[str, Any]]) -> list[dict[str, Any]]:
    """Group per-file action_required into a flat list of decision groups."""
    bucket: dict[str, list[str]] = {}
    options_map = {
        "ask_user_password_or_skip":
            ["password", "skip"],
        "ask_user_ocr_strategy":
            ["skip", "ocr_tesseract", "vlm_mlx", "claude_pagewise"],
        "ask_user_proceed_huge":
            ["skip", "proceed", "split"],
        "ask_user_skip_corrupt":
            ["skip"],
        "ask_user_skip_unsupported":
            ["skip"],
    }
    type_map = {
        "ask_user_password_or_skip": "encrypted",
        "ask_user_ocr_strategy": "scanned_pdf",
        "ask_user_proceed_huge": "huge_file",
        "ask_user_skip_corrupt": "corrupt",
        "ask_user_skip_unsupported": "unsupported_format",
    }
    defaults = {
        "ask_user_password_or_skip": "skip",
        "ask_user_ocr_strategy": "skip",
        "ask_user_proceed_huge": "skip",
        "ask_user_skip_corrupt": "skip",
        "ask_user_skip_unsupported": "skip",
    }
    for f in files:
        ar = f.get("action_required")
        if not ar:
            continue
        bucket.setdefault(ar, []).append(f["source_path"])

    out: list[dict[str, Any]] = []
    for action, paths in bucket.items():
        out.append({
            "type": type_map.get(action, action),
            "files": paths,
            "options": options_map.get(action, ["skip"]),
            "default": defaults.get(action, "skip"),
        })
    return out


def main() -> int:
    ap = argparse.ArgumentParser(description="Scan a document corpus and emit _scout.json.")
    ap.add_argument("input_dir", help="Path to the corpus root")
    ap.add_argument("kb_dir", help="Path to the kb/ output directory (will be created)")
    args = ap.parse_args()

    input_root = Path(args.input_dir).expanduser().resolve()
    kb_root = Path(args.kb_dir).expanduser().resolve()
    if not input_root.is_dir():
        log(f"input_dir does not exist or is not a directory: {input_root}")
        return 2

    kb_root.mkdir(parents=True, exist_ok=True)
    (kb_root / "docs").mkdir(exist_ok=True)
    (kb_root / "_logs").mkdir(exist_ok=True)

    t0 = time.time()
    paths, escapes = walk_corpus(input_root)
    log(f"Scanning {len(paths)} files under {input_root}")
    if escapes:
        log(f"Refused {len(escapes)} symlink escapes (see skipped_at_scout)")

    files: list[dict[str, Any]] = []
    skipped: list[dict[str, Any]] = list(escapes)
    for idx, p in enumerate(paths, start=1):
        try:
            info = scan_file(idx, p, input_root)
            files.append(info)
        except Exception as e:
            skipped.append({
                "source_path": str(p.relative_to(input_root)),
                "reason": f"scout failed: {str(e)[:200]}",
            })

    total_size = sum(f["size_bytes"] for f in files)
    total_tokens = sum((f.get("estimated_tokens") or 0) for f in files)
    elapsed = time.time() - t0
    # Cheap per-file extraction estimate (very rough — agent will refine).
    est_seconds = sum(
        2 if f.get("source_type") == "pdf" else 1
        for f in files
        if f.get("extraction_strategy") not in ("not_in_mvp", "skip",
                                                "needs_ocr_or_vlm", "needs_password")
    )

    scout: dict[str, Any] = {
        "schema_version": SCHEMA_VERSION,
        "scout_tool": _scout_tool_id(),
        "scanned_at": datetime.now(timezone.utc).strftime("%Y-%m-%dT%H:%M:%SZ"),
        "input_root": str(input_root),
        "kb_root": str(kb_root),
        "corpus": {
            "total_files": len(files),
            "total_size_bytes": total_size,
            "estimated_tokens": total_tokens,
            "estimated_extraction_seconds": est_seconds,
            "scout_elapsed_seconds": round(elapsed, 2),
        },
        "files": files,
        "skipped_at_scout": skipped,
        "user_decisions_needed": aggregate_decisions(files),
    }

    out_path = kb_root / "_scout.json"
    out_path.write_text(json.dumps(scout, ensure_ascii=False, indent=2), encoding="utf-8")
    log(f"Wrote {out_path}")

    # Human summary on stdout.
    types: dict[str, int] = {}
    strategies: dict[str, int] = {}
    for f in files:
        types[f["source_type"]] = types.get(f["source_type"], 0) + 1
        strategies[f["extraction_strategy"]] = strategies.get(f["extraction_strategy"], 0) + 1
    print(f"doc2kb scout: {len(files)} files, ~{total_tokens:,} tokens, {elapsed:.1f}s")
    print("  by type:      " + ", ".join(f"{k}={v}" for k, v in sorted(types.items())))
    print("  by strategy:  " + ", ".join(f"{k}={v}" for k, v in sorted(strategies.items())))
    if scout["user_decisions_needed"]:
        print(f"  USER DECISIONS NEEDED: {len(scout['user_decisions_needed'])} groups")
        for d in scout["user_decisions_needed"]:
            print(f"    - {d['type']}: {len(d['files'])} file(s); options={d['options']}; default={d['default']}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
