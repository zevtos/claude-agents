#!/usr/bin/env python3
"""
extract_pptx.py — Phase 4 extractor for PPTX files. Critically, this
preserves speaker notes — they often contain more semantic content than the
slides themselves, and most off-the-shelf converters drop them.

CLI:
    extract_pptx.py <input_pptx> <output_md>
                    [--doc-id doc-NNN]
                    [--source-rel relative/path/in/corpus.pptx]

Output layout (per slide):

    ## Slide N: <best-effort title>

    <body text from text frames, in shape order>

    | header | cells |
    | ---    | ---   |
    | ...    | ...   |

    ### Notes

    <speaker notes>

    ---  (slide separator)

The body intentionally interleaves shape contents in document order rather
than re-ordering by position — readers care about semantics, not layout.
"""
from __future__ import annotations

import argparse
import re
import sys
from pathlib import Path

from _common import (  # noqa: E402
    clean_whitespace,
    count_tokens,
    emit_failure,
    emit_success,
    sanitize_heading,
    sha256_of,
    today_iso,
    tool_version_string,
    validate_source_rel,
    write_md,
)


EXTRACTOR_NAME = "python-pptx"


def _try_import():
    try:
        from pptx import Presentation  # type: ignore
        from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore
        from pptx.enum.text import PP_ALIGN  # noqa: F401
        return Presentation, MSO_SHAPE_TYPE
    except Exception as e:
        emit_failure(f"python-pptx unavailable: {e}")
        sys.exit(1)


# Placeholder types that contain layout chrome (slide number, header, footer,
# date). These are noise for LLM consumption — skip them.
# Values from pptx.enum.text.PP_PLACEHOLDER:
#   13 SLIDE_NUMBER, 14 HEADER, 15 FOOTER, 16 DATE
_SKIP_PLACEHOLDER_TYPES = {13, 14, 15, 16}


def _is_chrome_placeholder(shape) -> bool:
    try:
        if not shape.is_placeholder:
            return False
        return int(shape.placeholder_format.type) in _SKIP_PLACEHOLDER_TYPES
    except Exception:
        return False


def _shape_text(shape) -> str:
    if not getattr(shape, "has_text_frame", False):
        return ""
    try:
        tf = shape.text_frame
    except Exception:
        return ""
    parts = []
    for para in tf.paragraphs:
        line = "".join(run.text for run in para.runs)
        if not line:
            # Sometimes runs are empty but para.text has content (auto-numbered, etc.)
            line = para.text or ""
        parts.append(line)
    return "\n".join(parts).rstrip()


def _table_md(shape) -> str:
    """Renders a pptx table as a Markdown pipe-table. Cells are joined with
    a space when they contain multiple paragraphs to avoid breaking the
    pipe-table format."""
    table = shape.table
    rows = []
    for row in table.rows:
        cells = []
        for cell in row.cells:
            tf = cell.text_frame
            txt = " ".join(p.text for p in tf.paragraphs if p.text).replace("\n", " ").strip()
            # Markdown pipes inside cells break the table; escape them.
            txt = txt.replace("|", "\\|")
            cells.append(txt)
        rows.append(cells)
    if not rows:
        return ""
    ncols = max(len(r) for r in rows)
    rows = [r + [""] * (ncols - len(r)) for r in rows]
    sep = ["---"] * ncols
    lines = ["| " + " | ".join(rows[0]) + " |",
             "| " + " | ".join(sep) + " |"]
    for r in rows[1:]:
        lines.append("| " + " | ".join(r) + " |")
    return "\n".join(lines)


def _best_title(slide, MSO_SHAPE_TYPE) -> str | None:
    # 1) explicit title placeholder
    for shape in slide.shapes:
        try:
            if shape.is_placeholder and shape.placeholder_format.idx == 0:
                t = _shape_text(shape).strip().split("\n", 1)[0]
                if t:
                    return t[:200]
        except Exception:
            pass
    # 2) first text shape's first line, if it looks like a heading
    for shape in slide.shapes:
        try:
            if shape.has_text_frame:
                t = _shape_text(shape).strip().split("\n", 1)[0]
                if t and len(t) < 200:
                    return t
        except Exception:
            pass
    return None


def extract(input_pptx: Path) -> tuple[str, dict, list[str]]:
    Presentation, MSO_SHAPE_TYPE = _try_import()
    warnings: list[str] = []
    extras: dict = {
        "slides": 0,
        "has_notes": False,
        "notes_chars": 0,
        "inline_images": 0,
        "has_tables": False,
        "has_charts": False,
    }

    try:
        prs = Presentation(str(input_pptx))
    except Exception as e:
        raise RuntimeError(f"failed to open pptx: {e}") from e

    extras["slides"] = len(prs.slides)
    pieces: list[str] = []
    n_images = 0
    n_tables = 0
    n_charts = 0
    total_notes = 0

    for idx, slide in enumerate(prs.slides, start=1):
        title = _best_title(slide, MSO_SHAPE_TYPE)
        if title:
            title = sanitize_heading(title)
        header = f"## Slide {idx}"
        if title:
            header += f": {title}"
        slide_parts: list[str] = [header, ""]

        title_shape_id = None
        # Find the title shape id so we don't print it twice.
        for shape in slide.shapes:
            try:
                if shape.is_placeholder and shape.placeholder_format.idx == 0:
                    title_shape_id = shape.shape_id
                    break
            except Exception:
                pass

        body_text_pieces: list[str] = []
        table_pieces: list[str] = []

        for shape in slide.shapes:
            try:
                stype = shape.shape_type
            except Exception:
                stype = None

            if _is_chrome_placeholder(shape):
                continue

            # Pictures: count only, no embedding.
            if stype == MSO_SHAPE_TYPE.PICTURE:
                n_images += 1
                continue

            # Tables.
            if getattr(shape, "has_table", False):
                n_tables += 1
                md = _table_md(shape)
                if md:
                    table_pieces.append(md)
                continue

            # Charts — we can't render them as markdown; record presence.
            if stype == MSO_SHAPE_TYPE.CHART or getattr(shape, "has_chart", False):
                n_charts += 1
                body_text_pieces.append("*(chart)*")
                continue

            # Text frames.
            if getattr(shape, "has_text_frame", False):
                if title_shape_id is not None and shape.shape_id == title_shape_id:
                    # Already used as slide header; skip duplicate.
                    continue
                txt = _shape_text(shape)
                if txt.strip():
                    body_text_pieces.append(txt.strip())
                continue

            # Grouped shapes — recurse one level to catch text inside groups.
            if stype == MSO_SHAPE_TYPE.GROUP:
                try:
                    for gsub in shape.shapes:
                        if getattr(gsub, "has_text_frame", False):
                            t = _shape_text(gsub).strip()
                            if t:
                                body_text_pieces.append(t)
                        elif gsub.shape_type == MSO_SHAPE_TYPE.PICTURE:
                            n_images += 1
                except Exception:
                    pass

        if body_text_pieces:
            slide_parts.append("\n\n".join(body_text_pieces))
        for tbl in table_pieces:
            slide_parts.append("")
            slide_parts.append(tbl)

        # Speaker notes.
        notes_text = ""
        if slide.has_notes_slide:
            try:
                notes_text = (slide.notes_slide.notes_text_frame.text or "").strip()
            except Exception:
                notes_text = ""
        if notes_text:
            extras["has_notes"] = True
            total_notes += len(notes_text)
            slide_parts.append("")
            slide_parts.append("### Notes")
            slide_parts.append("")
            slide_parts.append(notes_text)

        pieces.append("\n".join(slide_parts).rstrip())

    body = "\n\n---\n\n".join(pieces).strip() + "\n"
    body = clean_whitespace(body)

    extras["notes_chars"] = total_notes
    extras["inline_images"] = n_images
    extras["has_tables"] = n_tables > 0
    extras["has_charts"] = n_charts > 0
    if n_charts:
        warnings.append(f"{n_charts} chart(s) skipped (rendered as *(chart)* placeholder)")

    # Collect slide titles only, stripping the "Slide N: " chrome prefix so
    # `headings` is consistent with how other extractors record headings
    # (the actual title text, not the structural anchor).
    headings = []
    for line in body.split("\n"):
        if line.startswith("## ") and len(line) < 200:
            heading = line[3:]
            # "Slide N: Title" → "Title"; bare "Slide N" → skip (no title).
            m = re.match(r"^Slide \d+(?::\s*(.+))?$", heading)
            if m:
                title = (m.group(1) or "").strip()
                if not title:
                    continue
                heading = title
            sanitized = sanitize_heading(heading)
            if sanitized:
                headings.append(sanitized)
            if len(headings) >= 15:
                break
    extras["headings"] = headings

    return body, extras, warnings


def main() -> int:
    ap = argparse.ArgumentParser(description="Extract PPTX → Markdown via python-pptx.")
    ap.add_argument("input_pptx")
    ap.add_argument("output_md")
    ap.add_argument("--doc-id", default="doc-000")
    ap.add_argument("--source-rel", default=None)
    args = ap.parse_args()

    in_path = Path(args.input_pptx).expanduser().resolve()
    out_path = Path(args.output_md).expanduser().resolve()
    source_rel = args.source_rel or in_path.name
    try:
        source_rel = validate_source_rel(source_rel)
    except ValueError as e:
        emit_failure(f"invalid --source-rel: {e}")
        return 1

    if not in_path.is_file():
        emit_failure(f"input not found: {in_path}")
        return 1

    try:
        body, extras, warnings = extract(in_path)
    except Exception as e:
        emit_failure(f"extraction failed: {e}", extra={"input": str(in_path)})
        return 1

    fm = {
        "id": args.doc_id,
        "source": source_rel,
        "source_type": "pptx",
        "source_sha256": sha256_of(in_path),
        "extraction_method": f"{EXTRACTOR_NAME}@{tool_version_string()}",
        "extraction_date": today_iso(),
        "slides": extras.get("slides", 0),
        "has_notes": extras.get("has_notes", False),
        "notes_chars": extras.get("notes_chars", 0),
        "inline_images": extras.get("inline_images", 0),
        "has_tables": extras.get("has_tables", False),
        "has_charts": extras.get("has_charts", False),
        "headings": extras.get("headings", []),
        "tokens_estimated": count_tokens(body),
        "warnings": warnings,
    }
    write_md(out_path, fm, body)
    emit_success(out_path, body, warnings, extra={
        "slides": extras.get("slides", 0),
        "has_notes": extras.get("has_notes", False),
    })
    return 0


if __name__ == "__main__":
    sys.exit(main())
